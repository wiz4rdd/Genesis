"""
Autonomous Deep Research Agent – Hackathon Deck
Built on top of the HTF3_0-Idea-template.pptx template.

Template structure (20 × 11.25 inches):
  shape[0]  – background freeform  (keep)
  shape[1]  – Hack to Future logo  (keep)
  shape[2]  – footer text box      (keep)
  shape[3+] – content text boxes   (replace / add new ones)

Slides we produce (all inherit template backgrounds):
  1  – Title / Cover
  2  – Team Intro
  3  – Problem Overview
  4  – Solution
  5  – Tech Stack
  6  – Architecture & Flow
  7  – Agent Flow (steps)
  8  – Differentiation / Why Us?
  9  – Challenges, Future Scope & Business Potential
  10 – Demo / UI Mockup
  11 – Closing
"""

import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (matches original script) ──────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT  = RGBColor(0xF0, 0xF8, 0xFF)
TEAL    = RGBColor(0x00, 0xB4, 0xD8)
LTBLUE  = RGBColor(0x90, 0xE0, 0xEF)
GRAY    = RGBColor(0x88, 0xAA, 0xBB)
DKGRAY  = RGBColor(0x1E, 0x2D, 0x3D)
RED_ACC = RGBColor(0xFF, 0x55, 0x55)

# Template is 20 × 11.25 inches
W = 20.0
H = 11.25

# ── Helpers ───────────────────────────────────────────────────────────────────

def _emu(inches):
    return Emu(int(inches * 914400))

def clear_content_shapes(slide):
    """Remove all shapes beyond the first 3 (bg, logo, footer)."""
    sp_tree = slide.shapes._spTree
    to_remove = list(slide.shapes)[3:]
    for shape in to_remove:
        sp_tree.remove(shape._element)

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        _emu(x), _emu(y), _emu(w), _emu(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_txt(slide, text, x, y, w, h, color, size,
            bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txb

def add_label(slide, label, x, y, w, h, bg_color, text_color, size, bold=True):
    """Pill-style label: coloured rect + white text."""
    add_rect(slide, x, y, w, h, bg_color)
    add_txt(slide, label, x + 0.1, y + h * 0.1, w - 0.2, h * 0.8,
            text_color, size, bold=bold, align=PP_ALIGN.CENTER)

def add_card(slide, x, y, w, h, title, body, title_color=TEAL,
             body_color=WHITE, bg=DKGRAY, accent=TEAL):
    """Semi-transparent card with teal left accent + title + body."""
    add_rect(slide, x, y, w, h, bg)
    add_rect(slide, x, y, 0.08, h, accent)   # left accent bar
    add_txt(slide, title, x + 0.22, y + 0.12, w - 0.35, 0.55,
            title_color, 13, bold=True)
    add_txt(slide, body, x + 0.22, y + 0.68, w - 0.35, h - 0.80,
            body_color, 10, wrap=True)

# ── Load template ──────────────────────────────────────────────────────────────
TEMPLATE = '/mnt/user-data/uploads/HTF3_0-Idea-template.pptx'
prs = Presentation(TEMPLATE)

# We'll use the 11 template content slides (slides 2–12, index 1–11).
# Slide 1 (guidelines) we skip / remove later.
slides = list(prs.slides)   # 12 slides total

# Helper: get a slide by 0-based index
def S(i):
    return slides[i]

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1  → keep as guidelines slide (we'll delete shape[3] text and replace)
# Actually per template instructions, slide 1 is the guidelines slide to delete.
# We SKIP it and use slides 2–12 (indices 1–11) for our 11 content slides.
# ─────────────────────────────────────────────────────────────────────────────

# ── SLIDE 2: Title / Cover  (template index 1) ───────────────────────────────
s = S(1)
clear_content_shapes(s)

# Project ID + title block
add_txt(s, "AI-01  ·  Deep Research & Agentic AI", 0.6, 2.4, 19.0, 0.7,
        TEAL, 17, italic=True, align=PP_ALIGN.CENTER)

add_txt(s, "AUTONOMOUS", 0.5, 3.2, 19.0, 1.6,
        WHITE, 72, bold=True, align=PP_ALIGN.CENTER)
add_txt(s, "DEEP RESEARCH AGENT", 0.5, 4.75, 19.0, 1.1,
        TEAL, 48, bold=True, align=PP_ALIGN.CENTER)

add_txt(s, "AI-powered multi-step reasoning  ·  Live web search  ·  Source-backed synthesis",
        0.5, 5.95, 19.0, 0.65, LTBLUE, 18, italic=True, align=PP_ALIGN.CENTER)

add_txt(s, "🔍  LangGraph  ·  Claude / Groq  ·  Tavily Search API",
        0.5, 7.0, 19.0, 0.6, GRAY, 16, align=PP_ALIGN.CENTER)

add_txt(s, "Hackathon Submission  ·  Hack to Future 3.0  ·  2025",
        0.5, 9.8, 19.0, 0.55, GRAY, 14, align=PP_ALIGN.CENTER)

# ── SLIDE 3: Team Intro  (template index 2) ───────────────────────────────────
s = S(2)
clear_content_shapes(s)

add_txt(s, "MEET THE TEAM", 0.6, 2.2, 19.0, 0.8,
        TEAL, 34, bold=True, align=PP_ALIGN.CENTER)

team = [
    ("👤 Team Lead", "Aarav Shah", "Full-stack & AI integration"),
    ("👤 ML Engineer", "Priya Nair", "LangGraph agent design, LLM fine-tuning"),
    ("👤 Backend Dev", "Rohan Mehta", "FastAPI, Redis, deployment pipeline"),
    ("👤 Frontend Dev", "Sneha Patel", "Streamlit UI, UX & citations rendering"),
]

for i, (role, name, skill) in enumerate(team):
    col = i % 2
    row = i // 2
    x = 0.7 + col * 9.5
    y = 3.4 + row * 2.8
    add_rect(s, x, y, 8.8, 2.4, DKGRAY)
    add_rect(s, x, y, 0.12, 2.4, TEAL)
    add_txt(s, role, x + 0.3, y + 0.15, 8.3, 0.55, TEAL, 13, bold=True)
    add_txt(s, name, x + 0.3, y + 0.72, 8.3, 0.55, WHITE, 16, bold=True)
    add_txt(s, skill, x + 0.3, y + 1.3, 8.3, 0.8, GRAY, 12)

add_txt(s, "Motivation: Frustrated by hours spent on research — we built the tool we wished existed.",
        0.7, 9.35, 18.6, 0.6, LTBLUE, 13, italic=True, align=PP_ALIGN.CENTER)

# ── SLIDE 4: Problem Overview  (template index 3) ─────────────────────────────
s = S(3)
clear_content_shapes(s)

add_txt(s, "THE PROBLEM", 0.6, 2.0, 19.0, 0.75,
        TEAL, 34, bold=True, align=PP_ALIGN.CENTER)

problems = [
    ("📚", "Information Overload",
     "Sifting through dozens of tabs to answer one complex question wastes hours every day."),
    ("⏱️", "Time-Consuming Research",
     "Manual research — reading, cross-referencing, summarising — is slow, error-prone and exhausting."),
    ("❌", "No Source Traceability",
     "LLM answers lack citations, making it impossible to verify accuracy or trust the output."),
    ("👥", "Everyone Is Affected",
     "Students, researchers, analysts & professionals all face this daily — the pain is universal."),
]

for i, (icon, title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = 0.7 + col * 9.5
    y = 3.15 + row * 2.75
    add_card(s, x, y, 8.8, 2.45, f"{icon}  {title}", desc)

add_txt(s, '"Existing tools answer questions — they don\'t research them."',
        0.5, 9.3, 19.0, 0.6, TEAL, 14, bold=True, italic=True, align=PP_ALIGN.CENTER)

# ── SLIDE 5: Proposed Solution  (template index 4) ────────────────────────────
s = S(4)
clear_content_shapes(s)

add_txt(s, "THE SOLUTION", 0.6, 2.0, 19.0, 0.75,
        TEAL, 34, bold=True, align=PP_ALIGN.CENTER)

# Left column: What it does
add_txt(s, "What It Does", 0.7, 3.0, 8.5, 0.6, TEAL, 18, bold=True)
items = [
    "→  Accepts any natural language research question",
    "→  Autonomously searches the live web (Tavily API)",
    "→  Gathers & reads multiple sources in parallel",
    "→  Synthesises findings with multi-step LLM reasoning",
    "→  Returns concise answer with clickable citations",
    "→  Remembers session history for follow-up queries",
]
for j, item in enumerate(items):
    add_txt(s, item, 0.7, 3.75 + j * 0.9, 8.5, 0.78, WHITE, 13)

# Right column: Key differentiators
add_txt(s, "Key Differentiators", 10.4, 3.0, 8.5, 0.6, TEAL, 18, bold=True)
diffs = [
    ("🌐", "Live Web Data",        "Not stale training data — real-time results"),
    ("🔗", "Source Citations",     "Every claim fully traceable to a URL"),
    ("🧠", "Multi-Step Reasoning", "LangGraph agentic loops, not single-shot Q&A"),
    ("💬", "Conversational UI",    "Streamlit / React chat with session memory"),
    ("🗂️", "Session Memory",       "Context preserved across multiple turns"),
]
for j, (icon, head, sub) in enumerate(diffs):
    y = 3.75 + j * 1.05
    add_rect(s, 10.4, y, 9.0, 0.9, DKGRAY)
    add_rect(s, 10.4, y, 0.08, 0.9, TEAL)
    add_txt(s, icon + "  " + head, 10.6, y + 0.05, 4.5, 0.42, WHITE, 13, bold=True)
    add_txt(s, sub, 10.6, y + 0.5, 8.5, 0.38, GRAY, 11)

# ── SLIDE 6: Tech Stack  (template index 5) ───────────────────────────────────
s = S(5)
clear_content_shapes(s)

add_txt(s, "TECH STACK & METHODOLOGY", 0.6, 2.0, 19.0, 0.75,
        TEAL, 34, bold=True, align=PP_ALIGN.CENTER)

cats = [
    ("🔧 Orchestration",  ["LangGraph — agentic state machine",   "LangChain — tool & LLM wrappers"],
     "Chosen for fine-grained control of reasoning loops"),
    ("🤖 LLM",            ["Claude Sonnet (Anthropic)",            "Groq – LLaMA 3.1 (fast inference)"],
     "Best-in-class reasoning + speed trade-off"),
    ("🌐 Web Search",     ["Tavily Search API (live web)",         "Top-N results with full snippets"],
     "Purpose-built for LLM agents; not scraping"),
    ("🖥️ Frontend",       ["Streamlit (rapid prototyping)",        "React + FastAPI (production)"],
     "Rapid hackathon UI → production upgrade path"),
    ("🗂️ Memory",         ["Redis — session key-value store",      "SQLite — lightweight persistence"],
     "Enables multi-turn conversation context"),
    ("☁️ Deploy",         ["Render / Hugging Face Spaces",         ".env + python-dotenv for secrets"],
     "One-click deploy, zero infra overhead"),
]

for i, (cat, points, reason) in enumerate(cats):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 6.45
    y = 3.1 + row * 3.4
    add_rect(s, x, y, 6.1, 3.1, DKGRAY)
    add_rect(s, x, y, 6.1, 0.1, TEAL)
    add_txt(s, cat, x + 0.2, y + 0.2, 5.7, 0.55, TEAL, 14, bold=True)
    for j, pt in enumerate(points):
        add_txt(s, f"• {pt}", x + 0.2, y + 0.85 + j * 0.6, 5.7, 0.55, WHITE, 11)
    add_txt(s, f"↳ {reason}", x + 0.2, y + 2.1, 5.7, 0.75, GRAY, 10, italic=True)

# ── SLIDE 7: Architecture & Flow  (template index 6) ─────────────────────────
s = S(6)
clear_content_shapes(s)

add_txt(s, "SYSTEM ARCHITECTURE", 0.6, 2.0, 19.0, 0.75,
        TEAL, 34, bold=True, align=PP_ALIGN.CENTER)

# Four arch boxes
arch = [
    ("👤", "User",      "Chat UI\n(Streamlit/React)",  0.6),
    ("⚙️",  "Agent",    "LangGraph\nOrchestrator",     5.3),
    ("🔍", "Search",   "Tavily\nSearch API",           10.0),
    ("🧠", "LLM",      "Claude / Groq\nSynthesiser",   14.7),
]
for i, (icon, label, sub, x) in enumerate(arch):
    y = 3.2
    add_rect(s, x, y, 4.0, 4.0, DKGRAY)
    add_rect(s, x, y, 4.0, 0.12, TEAL)
    add_txt(s, icon, x + 0.2, y + 0.4, 3.6, 1.1, WHITE, 30, align=PP_ALIGN.CENTER)
    add_txt(s, label, x + 0.1, y + 1.55, 3.8, 0.5, TEAL, 14, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, sub, x + 0.1, y + 2.1, 3.8, 0.9, GRAY, 11, align=PP_ALIGN.CENTER)
    if i < 3:
        add_txt(s, "→", x + 4.05, y + 1.55, 0.55, 0.6, TEAL, 22, bold=True, align=PP_ALIGN.CENTER)

# Memory note below agent
add_rect(s, 5.3, 7.55, 4.0, 1.5, RGBColor(0x15, 0x28, 0x3A))
add_rect(s, 5.3, 7.55, 4.0, 0.1, TEAL)
add_txt(s, "↑  Session Memory\nRedis / SQLite Checkpointer",
        5.35, 7.72, 3.9, 1.2, WHITE, 11, align=PP_ALIGN.CENTER)

add_txt(s, "Live Web Sources: news · docs · Wikipedia · academic papers",
        10.0, 7.55, 9.0, 0.8, GRAY, 12, italic=True, align=PP_ALIGN.CENTER)

# ── SLIDE 8: Agent Flow  (template index 7) ───────────────────────────────────
s = S(7)
clear_content_shapes(s)

add_txt(s, "AGENT FLOW", 0.6, 2.0, 19.0, 0.75,
        TEAL, 34, bold=True, align=PP_ALIGN.CENTER)

steps = [
    ("1", "User Input",    "Question submitted via chat UI"),
    ("2", "LLM Reasoning", "LLM decides: answer directly OR search web?"),
    ("3", "Tool Call",     "Tavily API fetches top live web results"),
    ("4", "Read & Extract","Agent parses snippets, extracts key facts"),
    ("5", "Synthesise",    "LLM composes answer with inline citations"),
    ("6", "Return",        "Response streamed back to user in real time"),
]

for i, (num, head, sub) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 6.35
    y = 3.2 + row * 3.0
    add_rect(s, x, y, 5.9, 2.6, DKGRAY)
    add_rect(s, x, y, 1.0, 2.6, TEAL)
    add_txt(s, num, x + 0.15, y + 0.75, 0.75, 1.0, WHITE, 28, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, head, x + 1.2, y + 0.25, 4.5, 0.6, WHITE, 15, bold=True)
    add_txt(s, sub, x + 1.2, y + 0.95, 4.5, 1.3, LTBLUE, 11)
    if col < 2:
        add_txt(s, "→", x + 5.95, y + 0.85, 0.4, 0.7, TEAL, 18, bold=True, align=PP_ALIGN.CENTER)

# ── SLIDE 9: Differentiation  (template index 8) ─────────────────────────────
s = S(8)
clear_content_shapes(s)

add_txt(s, "DIFFERENTIATION — WHY US?", 0.6, 2.0, 19.0, 0.75,
        TEAL, 34, bold=True, align=PP_ALIGN.CENTER)

# Comparison table header
add_rect(s, 0.6, 3.0, 5.5, 0.65, TEAL)
add_rect(s, 6.2, 3.0, 6.2, 0.65, DKGRAY)
add_rect(s, 12.5, 3.0, 6.9, 0.65, RGBColor(0x02, 0x76, 0x3A))
add_txt(s, "Feature", 0.8, 3.07, 5.2, 0.52, WHITE, 13, bold=True)
add_txt(s, "Existing Tools (ChatGPT/Perplexity)", 6.3, 3.07, 5.9, 0.52, WHITE, 12, bold=True)
add_txt(s, "Our Agent", 12.6, 3.07, 6.5, 0.52, WHITE, 14, bold=True)

rows = [
    ("Live web data",        "Sometimes / stale",       "Always — Tavily real-time"),
    ("Source citations",     "Partial or none",         "Every claim cited"),
    ("Multi-step reasoning", "Single-shot",             "LangGraph agentic loops"),
    ("Session memory",       "Limited",                 "Redis / SQLite persisted"),
    ("Open source / extend", "Closed API",              "Fully open, hackable"),
]
for j, (feat, old, ours) in enumerate(rows):
    y = 3.75 + j * 0.95
    bg = RGBColor(0x12, 0x20, 0x30) if j % 2 == 0 else DKGRAY
    add_rect(s, 0.6, y, 5.5, 0.85, bg)
    add_rect(s, 6.2, y, 6.2, 0.85, bg)
    add_rect(s, 12.5, y, 6.9, 0.85, bg)
    add_txt(s, feat, 0.8, y + 0.15, 5.2, 0.6, LTBLUE, 12, bold=True)
    add_txt(s, old, 6.35, y + 0.15, 5.85, 0.6, GRAY, 11)
    add_txt(s, ours, 12.65, y + 0.15, 6.5, 0.6, WHITE, 11)

# ── SLIDE 10: Challenges + Future Scope  (template index 9) ──────────────────
s = S(9)
clear_content_shapes(s)

add_txt(s, "CHALLENGES, RISKS & FUTURE SCOPE", 0.6, 2.0, 19.0, 0.75,
        TEAL, 34, bold=True, align=PP_ALIGN.CENTER)

# Left: Challenges & Risks
add_txt(s, "⚠️  Challenges & Risks", 0.6, 3.05, 8.8, 0.65, TEAL, 18, bold=True)
challenges = [
    ("API Rate Limits",    "Tavily + Anthropic limits under high load"),
    ("Hallucination Risk", "LLM may mis-attribute even with citations"),
    ("Cost at Scale",      "Token usage grows with multi-hop queries"),
    ("Data Freshness",     "Tavily returns may lag breaking news"),
]
for j, (head, sub) in enumerate(challenges):
    y = 3.85 + j * 1.4
    add_rect(s, 0.6, y, 8.8, 1.25, DKGRAY)
    add_rect(s, 0.6, y, 0.1, 1.25, RED_ACC)
    add_txt(s, head, 0.85, y + 0.1, 8.3, 0.45, WHITE, 13, bold=True)
    add_txt(s, sub, 0.85, y + 0.6, 8.3, 0.55, GRAY, 11)

# Right: Future Scope
add_txt(s, "🚀  Future Scope & Business Potential", 10.0, 3.05, 9.4, 0.65, TEAL, 18, bold=True)
future = [
    ("📄 PDF / Doc ingestion",    "Research own documents alongside web"),
    ("🎙️ Voice interface",         "Speech-to-text query entry"),
    ("📊 Export reports",          "Auto-generate Word/PDF research briefs"),
    ("🏢 Enterprise SaaS",         "Multi-tenant knowledge assistant platform"),
    ("📈 Academic tool",           "Automated literature review & citing"),
]
for j, (head, sub) in enumerate(future):
    y = 3.85 + j * 1.28
    add_rect(s, 10.0, y, 9.4, 1.12, DKGRAY)
    add_rect(s, 10.0, y, 0.1, 1.12, TEAL)
    add_txt(s, head, 10.25, y + 0.1, 8.9, 0.45, WHITE, 13, bold=True)
    add_txt(s, sub, 10.25, y + 0.58, 8.9, 0.45, GRAY, 11)

# ── SLIDE 11: Demo / UI Mockup  (template index 10) ──────────────────────────
s = S(10)
clear_content_shapes(s)

add_txt(s, "DEMO — UI MOCKUP", 0.6, 2.0, 19.0, 0.75,
        TEAL, 34, bold=True, align=PP_ALIGN.CENTER)

# Browser chrome
add_rect(s, 0.8, 3.0, 18.4, 7.0, DKGRAY)
add_rect(s, 0.8, 3.0, 18.4, 0.65, RGBColor(0x2A, 0x3D, 0x52))
# Traffic lights
for cx, col in [(1.25, RGBColor(0xFF,0x5F,0x57)), (1.75, RGBColor(0xFF,0xBD,0x2E)), (2.25, RGBColor(0x28,0xC8,0x40))]:
    dot = s.shapes.add_shape(9, _emu(cx), _emu(3.18), _emu(0.28), _emu(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
add_txt(s, "🔍  Deep Research Agent  — Autonomous AI-powered research tool",
        2.7, 3.07, 13.0, 0.45, GRAY, 11)

# Sidebar
add_rect(s, 0.8, 3.65, 3.0, 6.35, RGBColor(0x15, 0x28, 0x3A))
add_txt(s, "Sessions", 1.0, 3.75, 2.6, 0.45, TEAL, 12, bold=True)
for j, sess in enumerate(["📌 LLM Reasoning", "📌 Climate Data", "📌 Market Analysis"]):
    add_txt(s, sess, 1.0, 4.35 + j * 0.65, 2.6, 0.55, GRAY, 10)

# Chat area
add_rect(s, 3.85, 3.65, 15.35, 5.6, RGBColor(0x0D, 0x1B, 0x2A))

# User bubble
add_rect(s, 9.5, 3.85, 9.4, 0.95, RGBColor(0x02, 0x76, 0x9E))
add_txt(s, "What are the latest breakthroughs in LLM reasoning?",
        9.65, 3.93, 9.0, 0.75, WHITE, 11)

# Agent bubble
add_rect(s, 3.95, 5.0, 14.5, 3.0, RGBColor(0x1E, 0x2D, 0x3D))
add_txt(s, "🤖  Deep Research Agent", 4.15, 5.1, 5.0, 0.45, TEAL, 12, bold=True)
add_txt(s,
    "Based on live sources, here are the top findings:\n"
    "1. Chain-of-thought prompting shows 40% improvement on math tasks  [Source 1]\n"
    "2. OpenAI o3 uses reinforcement learning for multi-step reasoning  [Source 2]\n"
    "3. Google Gemini 2.0 introduces agent-native reasoning loops        [Source 3]",
    4.15, 5.65, 14.0, 2.1, WHITE, 11)

# Input bar
add_rect(s, 3.85, 8.25, 15.35, 0.75, RGBColor(0x1E, 0x2D, 0x3D))
add_txt(s, "Ask anything… (Ctrl+Enter to send)",
        4.05, 8.38, 12.5, 0.5, GRAY, 11, italic=True)
add_rect(s, 18.15, 8.26, 1.0, 0.73, TEAL)
add_txt(s, "→", 18.2, 8.3, 0.9, 0.6, WHITE, 16, bold=True, align=PP_ALIGN.CENTER)

# ── SLIDE 12: Closing  (template index 11) ────────────────────────────────────
s = S(11)
clear_content_shapes(s)

add_txt(s, "THANK YOU", 0.5, 2.8, 19.0, 1.8,
        TEAL, 80, bold=True, align=PP_ALIGN.CENTER)
add_txt(s, "Autonomous Deep Research Agent", 0.5, 4.65, 19.0, 0.9,
        WHITE, 30, align=PP_ALIGN.CENTER)
add_txt(s, "Live Demo  ·  GitHub Repo  ·  Q&A",
        0.5, 5.7, 19.0, 0.65, LTBLUE, 18, italic=True, align=PP_ALIGN.CENTER)

links = [
    ("🔗", "github.com/your-handle/deep-research-agent"),
    ("📧", "your.email@example.com"),
    ("🤖", "Built with LangGraph · Claude · Tavily Search API"),
]
for j, (icon, text) in enumerate(links):
    add_txt(s, f"{icon}  {text}", 2.0, 7.0 + j * 0.85, 16.0, 0.75,
            GRAY, 14, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# REMOVE SLIDE 1 (the guidelines slide — index 0)
# python-pptx doesn't have a remove_slide method; use XML manipulation.
# ─────────────────────────────────────────────────────────────────────────────
from pptx.oxml.ns import qn

def remove_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides_part = prs.slides
    # Remove from the slide ID list
    slide_id_elem = xml_slides[index]
    xml_slides.remove(slide_id_elem)
    # Drop the slide part relationship
    slide = slides[index]
    rId = prs.slides._sldIdLst  # already removed — use Part.related_parts
    # Simpler: iterate rels and drop the one pointing to this slide
    for rel in prs.slides.part.rels.values():
        if rel.reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide':
            if rel._target == slide.part:
                prs.slides.part.drop_rel(rel.rId)
                break

remove_slide(prs, 0)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = '/home/claude/Deep_Research_Agent_HTF.pptx'
prs.save(OUT)
print("Saved:", OUT)